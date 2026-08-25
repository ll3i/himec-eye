# -*- coding: utf-8 -*-
"""발표자료(PPTX) 생성 — 붙임4 '첨부 자료'용.

붙임4 본문과 같은 내용·같은 수치를 쓴다. 문서와 발표자료가 어긋나면 안 되므로
성능 수치는 reports/performance.json 에서 직접 읽는다.
"""
import sys, json, argparse
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "reports" / "시연화면"

INK = RGBColor(0x19, 0x1F, 0x28)
MUTE = RGBColor(0x8B, 0x95, 0xA1)
SUB = RGBColor(0x4E, 0x59, 0x68)
BLUE = RGBColor(0x31, 0x82, 0xF6)
RED = RGBColor(0xF0, 0x44, 0x52)
LINE = RGBColor(0xE5, 0xE8, 0xEB)
BG = RGBColor(0xF9, 0xFA, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

W, H = Inches(13.333), Inches(7.5)


def txbox(slide, x, y, w, h, text, size=18, bold=False, color=INK,
          align=PP_ALIGN.LEFT, spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def rect(slide, x, y, w, h, fill=BG, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.06
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def fullbg(slide, color):
    """전면 배경 — 둥근 사각형을 쓰면 모서리에 흰 삼각이 남는다."""
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, kicker, title):
    txbox(slide, Inches(0.9), Inches(0.55), Inches(11), Inches(0.3),
          kicker, size=13, bold=True, color=BLUE)
    txbox(slide, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.8),
          title, size=30, bold=True)


def footer(slide, n):
    txbox(slide, Inches(11.6), Inches(6.9), Inches(0.9), Inches(0.3),
          str(n), size=11, color=MUTE, align=PP_ALIGN.RIGHT)


def table(slide, x, y, w, rows, col_w, head=True, size=13, rh=Inches(0.42)):
    """간단한 표 — 도형 대신 텍스트 상자로 그려 폰트를 제어한다."""
    yy = y
    for ri, row in enumerate(rows):
        if ri == 0 and head:
            rect(slide, x, yy, w, rh, fill=RGBColor(0xF2, 0xF4, 0xF6))
        xx = x
        for ci, cell in enumerate(row):
            txbox(slide, xx + Inches(0.12), yy + Inches(0.08), col_w[ci], rh,
                  str(cell), size=size, bold=(ri == 0 and head),
                  color=SUB if ri == 0 and head else INK)
            xx += col_w[ci]
        if ri > 0:
            ln = slide.shapes.add_shape(1, x, yy + rh, w, Emu(9525))
            ln.fill.solid(); ln.fill.fore_color.rgb = LINE
            ln.line.fill.background(); ln.shadow.inherit = False
        yy += rh
    return yy


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(ROOT / "reports" / "발표자료.pptx"))
    ap.add_argument("--youtube", default="https://youtu.be/KwVEdUqCQhc")
    ap.add_argument("--github", default="")
    ap.add_argument("--name", default="", help="표지 성명")
    ap.add_argument("--org", default="", help="표지 소속")
    ap.add_argument("--email", default="", help="마지막 장 연락처")
    a = ap.parse_args()

    reps = {r["task"]: r for r in json.loads(
        (ROOT / "reports" / "performance.json").read_text(encoding="utf-8"))}
    SUBMITTED = [("safety", "안전관리"), ("defect", "시설·설비 관리"),
                 ("product", "생산 품질검사"), ("polarity", "2차전지 셀 선별"),
                 ("logistics", "물류·운송 관리"), ("parking", "주차 운영관리")]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    n = 0

    # 1 표지
    s = blank(prs)
    fullbg(s, RGBColor(0x0D, 0x11, 0x16))
    txbox(s, Inches(1.1), Inches(2.1), Inches(11), Inches(0.5),
          "제1회 HIMEC AI 활용 아이디어 공모전  ·  ② 시공·품질관리",
          size=15, bold=True, color=RGBColor(0x5C, 0xA3, 0xDD))
    txbox(s, Inches(1.1), Inches(2.75), Inches(11), Inches(1.3),
          "HIMEC EYE — 직무별 비전 AI 관제", size=44, bold=True, color=WHITE)
    txbox(s, Inches(1.1), Inches(4.0), Inches(11), Inches(0.9),
          "현장 사진 1장에서 안전·품질 지적사항과 근거 조항·조치까지",
          size=21, color=RGBColor(0xB0, 0xB8, 0xC1))
    txbox(s, Inches(1.1), Inches(5.6), Inches(11), Inches(0.6),
          f"{a.name}  ·  {a.org}", size=15, color=RGBColor(0x8B, 0x95, 0xA1))
    n += 1

    # 2 제안 배경
    s = blank(prs); n += 1
    header(s, "BACKGROUND", "현장은 넓고, 판독은 사람 눈에 묶여 있습니다")
    txbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2),
          "삼구아이앤씨는 미화·보안·생산도급·물류·주차 운영처럼 사람이 직접 수행하는 현장 서비스를\n"
          "여러 사업장에서 동시에 운영합니다. 한 사업장 안에서도 생산 라인, 자재 야드, 기계실,\n"
          "주차 구획이 각각 다른 위험과 다른 품질 기준을 갖습니다.", size=17, color=SUB)
    box = [("판독이 사람에 묶임", "안전 점검자는 보호구를, 설비 점검자는 배관을 봅니다.\n같은 사진에 둘 다 있어도 목적 외는 지나칩니다."),
           ("기준의 속인화", "같은 부식·균열 사진도 경력에 따라\n경미/보통/중대 판정이 갈립니다."),
           ("사후 수기 기록", "찍고 → 정리하고 → 글로 옮기고 → 조항 찾기.\n사진이 늘수록 뒷단이 병목입니다."),
           ("직무별 기준 상이", "안전은 놓치면 사고, 품질은 오탐이면 라인 정지.\n한 기준을 전 직무에 못 씁니다.")]
    x = Inches(0.9)
    for t, d in box:
        rect(s, x, Inches(3.6), Inches(2.85), Inches(2.2))
        txbox(s, x + Inches(0.25), Inches(3.85), Inches(2.4), Inches(0.4), t, size=15, bold=True)
        txbox(s, x + Inches(0.25), Inches(4.35), Inches(2.4), Inches(1.4), d, size=12, color=MUTE)
        x += Inches(3.0)
    footer(s, n)

    # 3 제안
    s = blank(prs); n += 1
    header(s, "PROPOSAL", "만능 모델 하나가 아니라, 직무별 에이전트")
    txbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.9),
          "2차전지 셀의 극성 반전과 지게차 협착 위험은 찍는 대상도, 판정 기준도, 근거 법규도 다릅니다.\n"
          "직무마다 전용 모델과 전용 규칙을 가진 에이전트를 두고, 오케스트레이터가 위험도 순으로 합칩니다.",
          size=17, color=SUB)
    steps = [("1단계  지각", "직무별 객체탐지 모델", "분류가 아니라 탐지.\n'어디에 무엇이' 나와야\n위치 표시와 관계 판정이 됩니다."),
             ("2단계  판정", "결정론적 룰엔진 21개", "법규 판정에 LLM 미사용.\n조항을 잘못 인용한 문서는\n문서로서 가치를 잃습니다."),
             ("3단계  산출", "지적사항 · 준수율 · 이력", "규정번호·위험도·근거 조항·\n조치사항·좌표를 함께 남깁니다.")]
    x = Inches(0.9)
    for t, sub, d in steps:
        rect(s, x, Inches(3.5), Inches(3.75), Inches(2.5), fill=WHITE, line=LINE)
        txbox(s, x + Inches(0.3), Inches(3.75), Inches(3.2), Inches(0.35), t, size=13, bold=True, color=BLUE)
        txbox(s, x + Inches(0.3), Inches(4.15), Inches(3.2), Inches(0.4), sub, size=17, bold=True)
        txbox(s, x + Inches(0.3), Inches(4.7), Inches(3.2), Inches(1.2), d, size=12.5, color=MUTE)
        x += Inches(3.95)
    footer(s, n)

    # 4 6개 직무
    s = blank(prs); n += 1
    header(s, "SCOPE", "검증한 6개 직무")
    rows = [["직무", "검출 대상", "대표 판정 규칙"],
            ["안전관리", "작업자 · 보호구 착용/미착용 · 연기", "S-01 안전모 미착용 / S-05 연기 감지"],
            ["시설·설비 관리", "부식 · 균열 · 케이블 피복 · 용접 · 벽체", "D-02 균열 / D-03 케이블 피복 손상"],
            ["생산 품질검사", "냉납 · 조립 불량 · 단락 · 표면 결함", "P-02 조립 불량 / P-04 단락"],
            ["2차전지 셀 선별", "셀 단위 극성 방향", "X-01 역극 셀 검출"],
            ["물류·운송 관리", "지게차 · 작업자 · 섀시 적재 상태", "L-01 지게차 협착 위험"],
            ["주차 운영관리", "차종별 점유(승용·버스·트럭)", "K-01 점유 현황 집계"]]
    table(s, Inches(0.9), Inches(2.2), Inches(11.5), rows,
          [Inches(2.6), Inches(4.6), Inches(4.3)], rh=Inches(0.55))
    txbox(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.4),
          "판정 규칙 21개 — 각각 근거 조항 · 위험도 4단계 · 조치사항을 함께 갖습니다.",
          size=13, color=MUTE)
    footer(s, n)

    # 5 차별점 1 — 임계값 정책
    s = blank(prs); n += 1
    header(s, "KEY 1", "오탐과 미탐의 비용은 항목마다 다릅니다")
    txbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.6),
          "하나의 신뢰도 기준을 전 항목에 쓰지 않습니다. 항목마다 임계값을 다르게 두는 이유입니다.",
          size=17, color=SUB)
    rows = [["항목", "임계값", "방향", "이유"],
            ["안전모 미착용", "0.25", "재현율 우선", "놓치면 사고로 이어집니다"],
            ["연기 · 화재", "0.55", "정밀도 우선", "오경보를 남발하면 현장이 경보를 무시합니다"],
            ["역극 셀", "0.25", "재현율 최우선", "유출 시 클레임·안전사고"],
            ["정렬 정상 셀", "0.40", "보수적 확정", "확실할 때만 정상으로 확정합니다"],
            ["상하차 진행", "0.45", "정밀도 우선", "상태 판정이라 오탐이 잦으면 기록이 흐려집니다"]]
    table(s, Inches(0.9), Inches(2.9), Inches(11.5), rows,
          [Inches(2.8), Inches(1.4), Inches(2.2), Inches(5.1)], rh=Inches(0.52))
    txbox(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.5),
          "역극 셀 0.25 / 정상 셀 0.40 — 같은 공정 안에서도 방향이 반대인 정책을 씁니다.",
          size=13, bold=True, color=BLUE)
    footer(s, n)

    # 6 차별점 2 — 관계 추론 (스크린샷)
    s = blank(prs); n += 1
    header(s, "KEY 2", "각각 검출하는 것만으로는 나오지 않는 판정")
    shot = SHOTS / "1_관제화면_지게차협착위험경보.png"
    if shot.exists():
        s.shapes.add_picture(str(shot), Inches(0.9), Inches(2.05), width=Inches(7.4))
    txbox(s, Inches(8.6), Inches(2.1), Inches(3.9), Inches(3.5),
          "지게차 박스를 좌우 35% · 상하 17.5%\n확장한 영역에 작업자 박스가\n35% 이상 겹치면 협착 위험.\n\n"
          "지게차와 작업자를 각각 검출하는\n것만으로는 나오지 않습니다.\n두 객체의 '관계'가 판정 입력입니다.\n\n"
          "보호구도 같은 방식입니다.\n작업자는 있는데 머리 구간에\n안전모가 없으면 미착용으로 추론합니다.",
          size=14, color=SUB)
    txbox(s, Inches(8.6), Inches(5.7), Inches(3.9), Inches(0.8),
          "미착용 학습 데이터가 착용 대비 6.3:1로\n부족해, 모델 대신 규칙으로 보완했습니다.",
          size=12, color=MUTE)
    footer(s, n)

    # 7 시연 화면
    s = blank(prs); n += 1
    header(s, "DEMO", "심사위원이 직접 조작할 수 있습니다")
    pics = [("3_판독결과_직무별신뢰도와임계값.png", "판독 결과 — 신뢰도와 임계값을 함께 표시"),
            ("4_판정작동_클래스별임계값과채택여부.png", "판정 작동 — 채택/미채택이 그대로 드러남")]
    x = Inches(0.9)
    for f, cap in pics:
        p = SHOTS / f
        if p.exists():
            s.shapes.add_picture(str(p), x, Inches(2.2), width=Inches(5.6))
        txbox(s, x, Inches(5.1), Inches(5.6), Inches(0.4), cap, size=13, color=MUTE)
        x += Inches(5.95)
    txbox(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9),
          "https://himec-eye-panel.vercel.app\n"
          "관제 화면 · 현장 재현 영상 · 판독 결과 · 판정 작동 · 분석 리포트 (화면 5종)",
          size=15, bold=True, color=BLUE)
    footer(s, n)

    # 8 실증 성능
    s = blank(prs); n += 1
    header(s, "EVIDENCE", "아이디어가 아니라 실제로 도는 프로토타입")
    rows = [["직무", "mAP@50", "Precision", "Recall", "학습 epoch"]]
    for t, ko in SUBMITTED:
        r = reps.get(t)
        if r:
            o = r["overall"]
            rows.append([ko, f"{o['map50']*100:.1f}%", f"{o['precision']*100:.1f}%",
                         f"{o['recall']*100:.1f}%", str(r.get("epochs", "-"))])
    table(s, Inches(0.9), Inches(2.2), Inches(11.5), rows,
          [Inches(3.3), Inches(2.0), Inches(2.1), Inches(2.0), Inches(2.1)], rh=Inches(0.5))
    txbox(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(1.0),
          "· GPU 없이 CPU에서 제한된 epoch만 학습한 수치입니다. 수렴까지 학습하면 향상됩니다.\n"
          "· 2차전지 99.5%는 합성 데이터로 방법론을 검증한 값 — 도입 시 실제 셀 이미지로 재학습이 필요합니다.\n"
          "· 안전·생산이 낮은 것은 알고리즘이 아니라 데이터 문제입니다.",
          size=12.5, color=MUTE)
    footer(s, n)

    # 9 실제로 부딪힌 문제
    s = blank(prs); n += 1
    header(s, "LESSONS", "돌려보고 나서야 보인 것들")
    rows = [["문제", "원인", "해결", "결과"],
            ["연기 과탐지", "단일 임계값", "항목별 임계값 정책(연기 0.55)", "7건 → 2건"],
            ["하자 사진의 연기 오탐", "학습 분포 밖 이미지", "검측유형 라우팅", "3건 → 0건"],
            ["미착용 탐지 저조", "착용:미착용 = 6.3:1", "관계 추론 도입", "1건 → 5건"]]
    table(s, Inches(0.9), Inches(2.3), Inches(11.5), rows,
          [Inches(3.0), Inches(2.8), Inches(3.7), Inches(2.0)], rh=Inches(0.6))
    txbox(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.6),
          "두 번째 문제는 신뢰도로 거를 수 없었습니다 — 오탐 0.69 > 정탐 0.50.",
          size=14, bold=True, color=RED)
    footer(s, n)

    # 10 기대효과 · 로드맵
    s = blank(prs); n += 1
    header(s, "IMPACT", "기대효과와 도입 로드맵")
    items = [("판독 누락 감소", "직무 경계를 넘어 한 장에서 동시 판독"),
             ("기준의 표준화", "판정이 코드로 고정 — 경력에 따라 달라지지 않음"),
             ("문서 작성 부하 감소", "'쓰는 일'이 아니라 '확인하는 일'만"),
             ("확산 용이", "새 직무는 해당 모델과 규칙만 추가")]
    y = Inches(2.1)
    for t, d in items:
        rect(s, Inches(0.9), y, Inches(5.4), Inches(0.85), fill=WHITE, line=LINE)
        txbox(s, Inches(1.15), y + Inches(0.12), Inches(4.9), Inches(0.3), t, size=15, bold=True)
        txbox(s, Inches(1.15), y + Inches(0.46), Inches(4.9), Inches(0.3), d, size=12, color=MUTE)
        y += Inches(1.0)
    road = [("1단계  검증 (1~3개월)", "사내 현장 사진으로 재학습 · 오탐/미탐 실측 · 임계값 재조정"),
            ("2단계  현장 적용 (3~6개월)", "모바일 촬영 즉시 판독 · 점검 문서 연동 · 피드백 루프"),
            ("3단계  확산 (6~12개월)", "관리 시스템 연동 · 지적사항 DB · 협력사 품질 대시보드")]
    y = Inches(2.1)
    for t, d in road:
        txbox(s, Inches(6.9), y, Inches(5.5), Inches(0.35), t, size=15, bold=True, color=BLUE)
        txbox(s, Inches(6.9), y + Inches(0.38), Inches(5.5), Inches(0.6), d, size=12.5, color=SUB)
        y += Inches(1.3)
    txbox(s, Inches(6.9), Inches(6.0), Inches(5.5), Inches(0.8),
          "점진적 도입 — 초기에는 참고자료로만 제공하고,\n확인율·오탐률이 쌓인 뒤 자동화 범위를 넓힙니다.\n최종 판정 권한은 항상 책임자에게 있습니다.",
          size=12, color=MUTE)
    footer(s, n)

    # 11 링크
    s = blank(prs); n += 1
    fullbg(s, RGBColor(0x0D, 0x11, 0x16))
    txbox(s, Inches(1.1), Inches(1.6), Inches(11), Inches(0.6),
          "직접 확인해 보세요", size=32, bold=True, color=WHITE)
    links = [("온라인 시연", "https://himec-eye-panel.vercel.app",
              "관제 화면 · 현장 재현 영상 · 판독 결과 · 판정 작동 · 분석 리포트"),
             ("시연 영상", a.youtube, "관제 화면 36초 — 4개 현장, 판독부터 경보 확정까지")]
    if a.github:
        links.append(("소스코드", a.github, "데이터 수집 · 학습 · 판독 엔진 · 룰엔진 · 시연 페이지 빌더"))
    y = Inches(2.8)
    for t, url, d in links:
        txbox(s, Inches(1.1), y, Inches(2.2), Inches(0.4), t, size=15, bold=True,
              color=RGBColor(0x5C, 0xA3, 0xDD))
        txbox(s, Inches(3.3), y, Inches(9), Inches(0.4), url, size=17, bold=True, color=WHITE)
        txbox(s, Inches(3.3), y + Inches(0.42), Inches(9), Inches(0.4), d, size=12.5,
              color=RGBColor(0x8B, 0x95, 0xA1))
        y += Inches(1.15)
    txbox(s, Inches(1.1), Inches(6.5), Inches(11), Inches(0.4),
          f"{a.name}  ·  {a.org}  ·  {a.email}",
          size=12, color=RGBColor(0x8B, 0x95, 0xA1))

    prs.save(a.out)
    mb = Path(a.out).stat().st_size / 1e6
    print(f"{a.out}\n  슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장 · {mb:.1f} MB")


if __name__ == "__main__":
    main()
